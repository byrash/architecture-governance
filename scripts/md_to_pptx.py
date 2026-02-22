#!/usr/bin/env python3
"""Convert docs/lifecycle.md into a PowerPoint with native shapes from Mermaid diagrams.

Prerequisites (one-time):
    brew install mermaid-cli
    pip install python-pptx
"""

import re
import os
import json
import base64
import subprocess
import tempfile
from pathlib import Path
from html.parser import HTMLParser
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BG_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
TITLE_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
BODY_COLOR = RGBColor(0x44, 0x44, 0x44)

DIAGRAM_TOP = Inches(1.0)
DIAGRAM_MAX_W = SLIDE_WIDTH - Inches(0.6)
DIAGRAM_MAX_H = Inches(5.8)
DIAGRAM_LEFT_PAD = Inches(0.3)

PROJECT_ROOT = Path(__file__).resolve().parent.parent
SVG_NS = "{http://www.w3.org/2000/svg}"


# ---------------------------------------------------------------------------
# Markdown parsing (unchanged)
# ---------------------------------------------------------------------------

def parse_slides(md_path: Path) -> list[dict]:
    text = md_path.read_text()
    sections = re.split(r"\n---\n", text)
    slides = []
    for section in sections:
        title_m = re.search(r"^##\s+(.+)$", section, re.MULTILINE)
        if not title_m:
            continue
        raw_title = title_m.group(1).strip()
        title = re.sub(r"^Slide\s+\d+:\s*", "", raw_title)

        mermaid_m = re.search(r"```mermaid\n(.*?)```", section, re.DOTALL)
        mermaid_code = mermaid_m.group(1).strip() if mermaid_m else None

        legend_m = re.search(r'<p align="right">.*?</p>', section, re.DOTALL)
        legend_html = legend_m.group(0) if legend_m else None

        tp_m = re.search(
            r"### Talking Points\n(.*?)(?=### Speaker Notes|$)", section, re.DOTALL
        )
        talking_points = tp_m.group(1).strip() if tp_m else ""

        sn_m = re.search(r"### Speaker Notes\n(.*?)$", section, re.DOTALL)
        speaker_notes = sn_m.group(1).strip() if sn_m else ""

        slides.append(
            {
                "title": title,
                "mermaid": mermaid_code,
                "legend_html": legend_html,
                "talking_points": talking_points,
                "speaker_notes": speaker_notes,
            }
        )
    return slides


def parse_legend(html: str) -> list[tuple[str, str]]:
    pairs = []
    for m in re.finditer(r"background:(#[0-9A-Fa-f]{6})", html):
        color = m.group(1)
        after = html[m.end():]
        label_m = re.search(r"</span>\s*(.+?)(?:\s*<span|</sub>|$)", after)
        if label_m:
            label = re.sub(r"&nbsp;", " ", label_m.group(1)).strip().rstrip(";")
            pairs.append((color, label))
    return pairs


# ---------------------------------------------------------------------------
# Mermaid -> SVG rendering
# ---------------------------------------------------------------------------

def render_mermaid_svg(mermaid_code: str, out_svg: Path, config_path: Path) -> bool:
    with tempfile.NamedTemporaryFile(suffix=".mmd", mode="w", delete=False) as tmp:
        tmp.write(mermaid_code)
        tmp_path = tmp.name
    try:
        result = subprocess.run(
            ["mmdc", "-i", tmp_path, "-o", str(out_svg), "-c", str(config_path), "-b", "white"],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode != 0:
            print(f"  mmdc error: {result.stderr[:500]}")
            return False
        return out_svg.exists()
    finally:
        os.unlink(tmp_path)


# ---------------------------------------------------------------------------
# SVG parsing
# ---------------------------------------------------------------------------

class _TextExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self.lines = [""]

    def handle_starttag(self, tag, attrs):
        if tag == "br":
            self.lines.append("")

    def handle_data(self, data):
        self.lines[-1] += data.strip()

    def get_lines(self):
        return [l for l in self.lines if l]


def _html_text(html_str: str) -> list[str]:
    p = _TextExtractor()
    p.feed(html_str)
    return p.get_lines()


def _parse_transform(t: str | None):
    if not t:
        return 0.0, 0.0
    m = re.search(r"translate\(\s*([\d.\-]+)[,\s]+([\d.\-]+)\s*\)", t)
    return (float(m.group(1)), float(m.group(2))) if m else (0.0, 0.0)


def _parse_css(style: str | None) -> dict:
    if not style:
        return {}
    out = {}
    for part in style.split(";"):
        if ":" in part:
            k, v = part.split(":", 1)
            out[k.strip()] = v.replace("!important", "").strip()
    return out


def _fo_text(element) -> list[str]:
    for fo in element.iter(f"{SVG_NS}foreignObject"):
        return _html_text(ET.tostring(fo, encoding="unicode", method="html"))
    return []


def _decode_points(b64: str):
    try:
        return json.loads(base64.b64decode(b64))
    except Exception:
        return []


def _extract_svg_style(root) -> str:
    """Pull the full text from the <style> element."""
    for style_el in root.iter(f"{SVG_NS}style"):
        return style_el.text or ""
    return ""


def _css_class_prop(css_text: str, selector_fragment: str, prop: str, default: str) -> str:
    """Extract a CSS property from a class-based rule in the SVG <style> block."""
    pattern = re.compile(
        rf"\.{re.escape(selector_fragment)}\s*\{{[^}}]*{prop}\s*:\s*([^;}}]+)",
        re.IGNORECASE,
    )
    m = pattern.search(css_text)
    return m.group(1).strip() if m else default


def parse_svg(svg_path: Path) -> dict:
    tree = ET.parse(svg_path)
    root = tree.getroot()

    vb_str = root.get("viewBox", "0 0 100 100")
    vb = [float(x) for x in vb_str.split()]
    viewbox = {"x": vb[0], "y": vb[1], "w": vb[2], "h": vb[3]}

    svg_style = _extract_svg_style(root)
    cluster_fill = _css_class_prop(svg_style, "cluster rect", "fill", "#F3F4F6")
    cluster_stroke = _css_class_prop(svg_style, "cluster rect", "stroke", "#D1D5DB")

    clusters, nodes, edges, edge_labels = [], [], [], []

    def walk(el, px=0.0, py=0.0):
        tx, ty = _parse_transform(el.get("transform"))
        ax, ay = px + tx, py + ty

        cls_list = el.get("class", "").split()

        # Cluster (subgraph background)
        if "cluster" in cls_list and el.tag == f"{SVG_NS}g":
            rect = el.find(f"{SVG_NS}rect")
            if rect is not None:
                css = _parse_css(rect.get("style"))
                clusters.append(
                    {
                        "x": ax + float(rect.get("x", 0)),
                        "y": ay + float(rect.get("y", 0)),
                        "w": float(rect.get("width", 0)),
                        "h": float(rect.get("height", 0)),
                        "fill": css.get("fill") or cluster_fill,
                        "stroke": css.get("stroke") or cluster_stroke,
                        "label": " ".join(_fo_text(el)) or "",
                    }
                )

        # Node
        if "node" in cls_list and "nodes" not in cls_list and el.tag == f"{SVG_NS}g":
            rect = el.find(f"{SVG_NS}rect")
            if rect is not None:
                css = _parse_css(rect.get("style"))
                fill = css.get("fill", "#ECECFF")
                stroke = css.get("stroke", "#9370DB")

                text_color = "#333333"
                label_g = el.find(f"./{SVG_NS}g")
                if label_g is not None:
                    lc = _parse_css(label_g.get("style")).get("color")
                    if lc:
                        text_color = lc

                nodes.append(
                    {
                        "x": ax + float(rect.get("x", 0)),
                        "y": ay + float(rect.get("y", 0)),
                        "w": float(rect.get("width", 0)),
                        "h": float(rect.get("height", 0)),
                        "fill": fill,
                        "stroke": stroke,
                        "text": _fo_text(el),
                        "text_color": text_color,
                    }
                )

        # Edge path
        if el.tag == f"{SVG_NS}path" and el.get("data-edge") == "true":
            cls_str = el.get("class", "")
            if "edge-thickness-invisible" in cls_str:
                pass  # skip invisible links (~~~)
            else:
                pts_b64 = el.get("data-points", "")
                raw_pts = _decode_points(pts_b64) if pts_b64 else []
                if raw_pts:
                    abs_pts = [{"x": p["x"] + ax, "y": p["y"] + ay} for p in raw_pts]
                    has_arrow = "marker-end" in (el.get("marker-end") or el.attrib.get("marker-end", ""))
                    # check attribs more broadly
                    has_arrow = any("pointEnd" in v for v in el.attrib.values())
                    dashed = "edge-pattern-dashed" in cls_str or "edge-pattern-dotted" in cls_str
                    edges.append({"points": abs_pts, "arrow": has_arrow, "dashed": dashed})

        # Edge label
        if "edgeLabel" in cls_list and el.tag == f"{SVG_NS}g":
            label_g = el.find(f".//{SVG_NS}g[@class='label']")
            if label_g is not None:
                lx, ly = _parse_transform(label_g.get("transform"))
                text = " ".join(_fo_text(label_g)).strip()
                if text:
                    edge_labels.append({"x": ax + lx, "y": ay + ly, "text": text})

        for child in el:
            walk(child, ax, ay)

    walk(root)
    return {
        "viewbox": viewbox,
        "clusters": clusters,
        "nodes": nodes,
        "edges": edges,
        "edge_labels": edge_labels,
    }


# ---------------------------------------------------------------------------
# PowerPoint rendering helpers
# ---------------------------------------------------------------------------

def _hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    if len(h) == 3:
        h = h[0] * 2 + h[1] * 2 + h[2] * 2
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _css_color_to_hex(c: str) -> str:
    """Handle 'rgb(r, g, b)' or '#hex'."""
    c = c.strip()
    m = re.match(r"rgb\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)", c)
    if m:
        return "#{:02x}{:02x}{:02x}".format(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    return c


def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_arrow_end(connector):
    """Add triangle arrowhead to the end of a connector via oxml."""
    ln = connector._element.spPr.find(qn("a:ln"))
    if ln is None:
        from lxml import etree
        ln = etree.SubElement(connector._element.spPr, qn("a:ln"))
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        from lxml import etree
        tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")


# ---------------------------------------------------------------------------
# Build slides
# ---------------------------------------------------------------------------

def add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)
    txbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), SLIDE_WIDTH - Inches(2), Inches(2))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Architecture Governance"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Lifecycle & Governance Model"
    p2.font.size = Pt(24)
    p2.font.color.rgb = BODY_COLOR
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)


def add_content_slide(prs, title, svg_data, legend_items, talking_points, speaker_notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), SLIDE_WIDTH - Inches(1), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # Render native shapes from SVG data
    if svg_data:
        _render_shapes(slide, svg_data)

    # Legend
    if legend_items:
        legend_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.7), SLIDE_WIDTH - Inches(1), Inches(0.4)
        )
        tf_leg = legend_box.text_frame
        tf_leg.word_wrap = True
        p_leg = tf_leg.paragraphs[0]
        p_leg.alignment = PP_ALIGN.RIGHT
        for i, (color, label) in enumerate(legend_items):
            run_dot = p_leg.add_run()
            run_dot.text = "\u25A0 "
            run_dot.font.size = Pt(10)
            run_dot.font.color.rgb = _hex_to_rgb(color)
            run_lbl = p_leg.add_run()
            run_lbl.text = label + ("     " if i < len(legend_items) - 1 else "")
            run_lbl.font.size = Pt(10)
            run_lbl.font.color.rgb = BODY_COLOR

    # Speaker notes
    notes_tf = slide.notes_slide.notes_text_frame
    if talking_points:
        notes_tf.text = "TALKING POINTS:\n" + talking_points
    if speaker_notes:
        p_sn = notes_tf.add_paragraph()
        p_sn.text = "\nSPEAKER NOTES:\n" + speaker_notes


def _render_shapes(slide, svg_data):
    """Render SVG elements as native PowerPoint shapes on the slide."""
    vb = svg_data["viewbox"]
    svg_w, svg_h = vb["w"], vb["h"]
    if svg_w == 0 or svg_h == 0:
        return

    max_w = DIAGRAM_MAX_W
    max_h = DIAGRAM_MAX_H

    scale_x = max_w / svg_w
    scale_y = max_h / svg_h
    scale = min(scale_x, scale_y)

    rendered_w = int(svg_w * scale)
    rendered_h = int(svg_h * scale)
    offset_x = DIAGRAM_LEFT_PAD + (max_w - rendered_w) // 2
    offset_y = DIAGRAM_TOP + (max_h - rendered_h) // 2

    def sx(v):
        return int(offset_x + v * scale)

    def sy(v):
        return int(offset_y + v * scale)

    def sw(v):
        return int(v * scale)

    # 1. Cluster backgrounds (draw first so they're behind nodes)
    for c in svg_data["clusters"]:
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            sx(c["x"]), sy(c["y"]), sw(c["w"]), sw(c["h"]),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(_css_color_to_hex(c["fill"]))
        shape.line.color.rgb = _hex_to_rgb(_css_color_to_hex(c["stroke"]))
        shape.line.width = Pt(1)

        if c["label"]:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.auto_size = None
            p = tf.paragraphs[0]
            p.text = c["label"]
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].space_before = Pt(0)
            shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # 2. Edges (draw before nodes so nodes cover edge ends cleanly)
    for e in svg_data["edges"]:
        pts = e["points"]
        if len(pts) < 2:
            continue
        start, end = pts[0], pts[-1]
        cxn = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            sx(start["x"]), sy(start["y"]),
            sx(end["x"]), sy(end["y"]),
        )
        cxn.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        cxn.line.width = Pt(1.5)
        if e["dashed"]:
            cxn.line.dash_style = 2  # MSO_LINE_DASH_STYLE.DASH
        if e["arrow"]:
            _add_arrow_end(cxn)

    # 3. Edge labels
    for el in svg_data["edge_labels"]:
        lw, lh = Inches(1.5), Inches(0.35)
        lx = sx(el["x"]) - lw // 2
        ly = sy(el["y"]) - lh // 2
        tb = slide.shapes.add_textbox(lx, ly, lw, lh)
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.text = el["text"]
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.alignment = PP_ALIGN.CENTER

    # 4. Nodes (on top)
    for n in svg_data["nodes"]:
        shape = slide.shapes.add_shape(
            5,  # MSO_SHAPE.ROUNDED_RECTANGLE
            sx(n["x"]), sy(n["y"]), sw(n["w"]), sw(n["h"]),
        )
        fill_hex = _css_color_to_hex(n["fill"])
        stroke_hex = _css_color_to_hex(n["stroke"])
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
        shape.line.color.rgb = _hex_to_rgb(stroke_hex)
        shape.line.width = Pt(1)

        text_hex = _css_color_to_hex(n["text_color"])
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        for i, line in enumerate(n["text"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(10)
            p.font.color.rgb = _hex_to_rgb(text_hex)
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(0)
            p.space_after = Pt(0)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    md_path = PROJECT_ROOT / "docs" / "lifecycle.md"
    out_dir = PROJECT_ROOT / "output"
    svg_dir = out_dir / "svg"
    svg_dir.mkdir(parents=True, exist_ok=True)

    config_path = out_dir / "mermaid-config.json"
    config_path.write_text(json.dumps({
        "theme": "base",
        "themeVariables": {
            "clusterBkg": "#F3F4F6",
            "clusterBorder": "#D1D5DB",
        },
    }))

    print("Parsing lifecycle.md ...")
    slides = parse_slides(md_path)
    print(f"Found {len(slides)} slides")

    for i, s in enumerate(slides):
        if s["mermaid"]:
            svg_path = svg_dir / f"slide_{i + 1}.svg"
            print(f"Rendering Slide {i + 1}: {s['title']} ...")
            ok = render_mermaid_svg(s["mermaid"], svg_path, config_path)
            if ok:
                s["svg_data"] = parse_svg(svg_path)
            else:
                print(f"  WARNING: Failed to render slide {i + 1}")
                s["svg_data"] = None
        else:
            s["svg_data"] = None

    print("Building PowerPoint ...")
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    add_title_slide(prs)

    for s in slides:
        legend_items = parse_legend(s["legend_html"]) if s["legend_html"] else []
        add_content_slide(
            prs, s["title"], s["svg_data"], legend_items,
            s["talking_points"], s["speaker_notes"],
        )

    pptx_path = out_dir / "architecture-governance-lifecycle.pptx"
    prs.save(str(pptx_path))
    print(f"\nDone! PowerPoint saved to: {pptx_path}")


if __name__ == "__main__":
    main()
