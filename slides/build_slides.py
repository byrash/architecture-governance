#!/usr/bin/env python3
"""Build 2 executive architecture slides with native PowerPoint shapes from Mermaid.

Reuses the SVG-parsing and shape-rendering pipeline from scripts/md_to_pptx.py.
"""

import json
import os
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

import sys
PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT / "scripts"))
from md_to_pptx import (
    parse_svg,
    _render_shapes,
    set_slide_bg,
    _hex_to_rgb,
    SLIDE_WIDTH,
    SLIDE_HEIGHT,
    BG_COLOR,
    TITLE_COLOR,
    BODY_COLOR,
    DIAGRAM_TOP,
    DIAGRAM_MAX_W,
    DIAGRAM_MAX_H,
    DIAGRAM_LEFT_PAD,
)

SLIDES_DIR = Path(__file__).parent
PUPPETEER_CONFIG = SLIDES_DIR / "puppeteer-config.json"
MERMAID_CONFIG = SLIDES_DIR / "mermaid-config.json"

LEGEND_ITEMS = [
    ("#6c757d", "Input"),
    ("#2b7bba", "Deterministic"),
    ("#e8833a", "LLM-Powered"),
    ("#3a9d5c", "Output"),
]

SLIDES = [
    {
        "title": "Ingestion: Building the Knowledge Base",
        "mmd_file": "slide1-ingestion.mmd",
        "svg_file": "slide1-ingestion.svg",
        "bullets": [
            "Input: Confluence reference pages with embedded diagrams (Draw.io, SVG, PlantUML)",
            "Deterministic: Diagrams parsed to AST, structural rules extracted with zero LLM usage",
            "LLM used only twice: text rule extraction from prose + rule enrichment (both cached, run once per change)",
            "Per-page rules bubble up into consolidated _all.rules.md across all pages in each index",
            "Output: Knowledge base per index (_all.rules.md + rules-enriched.json with synonyms, patterns, hints)",
        ],
    },
    {
        "title": "Validation: Scoring a Target Page",
        "mmd_file": "slide2-validation.mmd",
        "svg_file": "slide2-validation.svg",
        "bullets": [
            "Target page downloaded from Confluence, diagrams parsed to AST (same pipeline as ingestion)",
            "Scored against 3 knowledge base indexes: Security (Weight: 40%), Patterns (Weight: 30%), Standards (Weight: 30%)",
            "80\u201390% of scoring is deterministic Python \u2014 stable, repeatable, no LLM variance",
            "LLM validation agents handle only 10\u201320% ambiguous rules (weak evidence or contradictions)",
            "Output: Governance report, HTML dashboard, Confluence comment (Pass \u2265 70, Warn 50\u201369, Fail < 50)",
        ],
    },
]


def render_mermaid_to_svg(mmd_path: Path, svg_path: Path) -> bool:
    cmd = ["mmdc", "-i", str(mmd_path), "-o", str(svg_path),
           "-p", str(PUPPETEER_CONFIG), "-c", str(MERMAID_CONFIG), "-b", "white"]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
    if result.returncode != 0:
        print(f"  mmdc error: {result.stderr[:500]}")
        return False
    return svg_path.exists()


def add_slide(prs, title, svg_data, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), SLIDE_WIDTH - Inches(1), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.LEFT

    # Native diagram shapes
    if svg_data:
        _render_shapes(slide, svg_data)

    # Bullet points
    bullet_top = Inches(5.2)
    bullet_box = slide.shapes.add_textbox(
        Inches(0.7), bullet_top, SLIDE_WIDTH - Inches(1.4), Inches(1.9))
    tf = bullet_box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(12)
        p.font.color.rgb = BODY_COLOR
        p.space_after = Pt(4)
        pf = p._pPr
        bu = pf.makeelement(qn("a:buChar"), {"char": "\u2022"})
        pf.append(bu)

    # Color legend at bottom
    legend_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.0), SLIDE_WIDTH - Inches(1), Inches(0.35))
    tf_leg = legend_box.text_frame
    tf_leg.word_wrap = True
    p_leg = tf_leg.paragraphs[0]
    p_leg.alignment = PP_ALIGN.RIGHT
    for i, (color, label) in enumerate(LEGEND_ITEMS):
        run_dot = p_leg.add_run()
        run_dot.text = "\u25A0 "
        run_dot.font.size = Pt(10)
        run_dot.font.color.rgb = _hex_to_rgb(color)
        run_lbl = p_leg.add_run()
        run_lbl.text = label + ("     " if i < len(LEGEND_ITEMS) - 1 else "")
        run_lbl.font.size = Pt(10)
        run_lbl.font.color.rgb = BODY_COLOR


def main():
    # Write configs
    PUPPETEER_CONFIG.write_text(json.dumps({
        "executablePath": "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "args": ["--no-sandbox"],
    }))
    MERMAID_CONFIG.write_text(json.dumps({
        "theme": "base",
        "themeVariables": {
            "clusterBkg": "#F3F4F6",
            "clusterBorder": "#D1D5DB",
        },
    }))

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for s in SLIDES:
        mmd_path = SLIDES_DIR / s["mmd_file"]
        svg_path = SLIDES_DIR / s["svg_file"]

        print(f"Rendering: {s['title']} ...")
        ok = render_mermaid_to_svg(mmd_path, svg_path)
        if not ok:
            print(f"  FAILED — skipping native shapes for this slide")
            svg_data = None
        else:
            svg_data = parse_svg(svg_path)
            print(f"  Parsed: {len(svg_data['nodes'])} nodes, {len(svg_data['edges'])} edges, "
                  f"{len(svg_data['clusters'])} clusters")

        add_slide(prs, s["title"], svg_data, s["bullets"])

    out_path = SLIDES_DIR / "architecture-governance.pptx"
    prs.save(str(out_path))
    print(f"\nDone! {out_path}")

    # Cleanup temp configs
    PUPPETEER_CONFIG.unlink(missing_ok=True)
    MERMAID_CONFIG.unlink(missing_ok=True)


if __name__ == "__main__":
    main()
